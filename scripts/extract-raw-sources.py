#!/usr/bin/env python3
"""raw 소스 자동 추출기: PDF, PPTX, DOCX, XLSX → 마크다운

markitdown(microsoft/markitdown)을 1차 엔진으로 사용하고,
pymupdf/python-pptx/openpyxl을 폴백으로 유지한다.

사용법:
  python3 scripts/extract-raw-sources.py [소스폴더] [출력폴더] [--type pdf|pptx|docx|xlsx|all]
  # markitdown venv 사용 시:
  .venv/bin/python scripts/extract-raw-sources.py raw/security-onsite-reports raw/extracted --type pptx

예시:
  .venv/bin/python scripts/extract-raw-sources.py raw/security-onsite-reports raw/extracted --type pptx
  .venv/bin/python scripts/extract-raw-sources.py raw/security-microsoft-documents raw/extracted --type all
  # 폴백 전용 (markitdown 없는 환경):
  python3 scripts/extract-raw-sources.py raw/security-onsite-reports raw/extracted --type pptx --no-markitdown
"""
import sys, os, argparse, json
from pathlib import Path
from datetime import datetime

# ─── markitdown (primary engine) ───
_markitdown_instance = None

def _get_markitdown():
    """markitdown 인스턴스를 싱글턴으로 반환. 없으면 None."""
    global _markitdown_instance
    if _markitdown_instance is False:
        return None
    if _markitdown_instance is None:
        try:
            from markitdown import MarkItDown
            _markitdown_instance = MarkItDown()
        except ImportError:
            _markitdown_instance = False
            return None
    return _markitdown_instance

def extract_with_markitdown(filepath):
    """markitdown으로 파일→마크다운 변환."""
    md = _get_markitdown()
    if md is None:
        return None
    try:
        result = md.convert(str(filepath))
        if result and result.text_content and result.text_content.strip():
            return result.text_content
    except Exception:
        pass
    return None

# ─── PDF 추출 (pymupdf - PyMuPDF) ───
def extract_pdf(filepath):
    """PDF에서 텍스트 추출. pymupdf(fitz) 사용."""
    try:
        import fitz
        doc = fitz.open(filepath)
        name = Path(filepath).stem
        lines = [
            f"# {name}",
            f"",
            f"- **소스**: `{filepath}`",
            f"- **페이지 수**: {len(doc)}",
            f"- **추출일**: {datetime.now().strftime('%Y-%m-%d')}",
            ""
        ]
        for i, page in enumerate(doc, 1):
            text = page.get_text("text").strip()
            if text:
                lines.append(f"## Page {i}")
                lines.append(text)
                lines.append("")
        doc.close()
        return "\n".join(lines)
    except Exception as e:
        return f"# {Path(filepath).stem}\n\nERROR: {e}\n"


# ─── PPTX 추출 (python-pptx) ───
def extract_pptx(filepath):
    """PPTX에서 슬라이드별 텍스트 추출."""
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        name = Path(filepath).stem
        lines = [
            f"# {name}",
            f"",
            f"- **소스**: `{filepath}`",
            f"- **슬라이드 수**: {len(prs.slides)}",
            f"- **추출일**: {datetime.now().strftime('%Y-%m-%d')}",
            ""
        ]
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells
                        )
                        if row_text.strip(" |"):
                            texts.append(f"| {row_text} |")
            if texts:
                lines.append(f"## Slide {i}")
                lines.append("\n".join(texts))
                lines.append("")
            # 발표자 노트
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes and notes != str(i):
                    lines.append(f"> **노트**: {notes[:500]}")
                    lines.append("")
        return "\n".join(lines)
    except Exception as e:
        return f"# {Path(filepath).stem}\n\nERROR: {e}\n"


# ─── DOCX 추출 (python-docx 없으면 zipfile + xml) ───
def extract_docx(filepath):
    """DOCX에서 텍스트 추출. python-docx 또는 zipfile fallback."""
    try:
        import zipfile, xml.etree.ElementTree as ET
        name = Path(filepath).stem
        
        with zipfile.ZipFile(filepath) as z:
            # document.xml에서 텍스트 추출
            doc_xml = z.read('word/document.xml')
            root = ET.fromstring(doc_xml)
            
            ns = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
            paragraphs = root.findall('.//w:p', ns)
            
            lines = [
                f"# {name}",
                f"",
                f"- **소스**: `{filepath}`",
                f"- **문단 수**: {len(paragraphs)}",
                f"- **추출일**: {datetime.now().strftime('%Y-%m-%d')}",
                ""
            ]
            
            for para in paragraphs:
                texts = []
                for run in para.findall('.//w:r', ns):
                    for t in run.findall('.//w:t', ns):
                        if t.text:
                            texts.append(t.text)
                full = ''.join(texts).strip()
                if full:
                    lines.append(full)
            
            return "\n\n".join(lines) if lines else f"# {name}\n\n(빈 문서)\n"
    except Exception as e:
        return f"# {Path(filepath).stem}\n\nERROR: {e}\n"


# ─── XLSX 추출 (openpyxl) ───
def extract_xlsx(filepath):
    """XLSX에서 시트별 테이블 추출."""
    try:
        from openpyxl import load_workbook
        name = Path(filepath).stem
        wb = load_workbook(filepath, read_only=True, data_only=True)
        
        lines = [
            f"# {name}",
            f"",
            f"- **소스**: `{filepath}`",
            f"- **시트 수**: {len(wb.sheetnames)}",
            f"- **추출일**: {datetime.now().strftime('%Y-%m-%d')}",
            ""
        ]
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                continue
            
            lines.append(f"## {sheet_name}")
            lines.append("")
            
            # 마크다운 테이블로 변환 (상위 100행)
            for i, row in enumerate(rows[:100]):
                cells = [str(c) if c is not None else "" for c in row]
                lines.append("| " + " | ".join(cells) + " |")
                if i == 0:
                    lines.append("| " + " | ".join(["---"] * len(cells)) + " |")
            
            if len(rows) > 100:
                lines.append(f"\n_... 총 {len(rows)}행 중 상위 100행만 표시_")
            lines.append("")
        
        wb.close()
        return "\n".join(lines)
    except Exception as e:
        return f"# {Path(filepath).stem}\n\nERROR: {e}\n"


# ─── 메인 ───
EXTRACTORS = {
    '.pdf': extract_pdf,
    '.pptx': extract_pptx,
    '.docx': extract_docx,
    '.xlsx': extract_xlsx,
}

EXT_MAP = {
    'pdf': ['.pdf'],
    'pptx': ['.pptx'],
    'docx': ['.docx'],
    'xlsx': ['.xlsx'],
    'all': ['.pdf', '.pptx', '.docx', '.xlsx'],
}


def main():
    parser = argparse.ArgumentParser(description="raw 소스 바이너리→마크다운 추출")
    parser.add_argument('source', help='소스 폴더 (예: raw/security-onsite-reports)')
    parser.add_argument('output', help='출력 폴더 (예: raw/extracted)')
    parser.add_argument('--type', default='all', choices=EXT_MAP.keys(),
                        help='추출할 파일 유형 (기본: all)')
    parser.add_argument('--limit', type=int, default=0,
                        help='최대 파일 수 (0=무제한)')
    parser.add_argument('--force', action='store_true',
                        help='기존 파일 덮어쓰기')
    parser.add_argument('--dry-run', action='store_true',
                        help='실제 추출 없이 대상 파일만 나열')
    parser.add_argument('--no-markitdown', action='store_true',
                        help='markitdown 사용 안 함 (폴백 추출기만 사용)')
    args = parser.parse_args()

    source = Path(args.source)
    output = Path(args.output)
    exts = EXT_MAP[args.type]

    if args.no_markitdown:
        global _markitdown_instance
        _markitdown_instance = False

    # markitdown 사용 가능 여부 확인
    use_markitdown = _get_markitdown() is not None
    engine_label = "markitdown" if use_markitdown else "fallback"
    print(f"엔진: {engine_label}")

    if not source.exists():
        print(f"ERROR: 소스 폴더 없음: {source}")
        sys.exit(1)

    # 대상 파일 수집
    files = []
    for ext in exts:
        files.extend(sorted(source.rglob(f'*{ext}')))
        files.extend(sorted(source.rglob(f'*{ext.upper()}')))
    
    # 중복 제거
    files = list(dict.fromkeys(files))
    
    if args.limit:
        files = files[:args.limit]

    print(f"대상: {len(files)}개 파일 ({', '.join(exts)})")

    if args.dry_run:
        for f in files[:30]:
            print(f"  {f}")
        if len(files) > 30:
            print(f"  ... +{len(files) - 30}개")
        return

    # 추출
    stats = {'ok': 0, 'skip': 0, 'error': 0}
    
    for i, filepath in enumerate(files, 1):
        rel = filepath.relative_to(source)
        out_path = output / rel.with_suffix('.md')

        if out_path.exists() and not args.force:
            stats['skip'] += 1
            continue

        ext = filepath.suffix.lower()
        extractor = EXTRACTORS.get(ext)
        if not extractor:
            continue

        try:
            # 1차: markitdown
            content = extract_with_markitdown(str(filepath))
            # 2차: 폴백 추출기
            if content is None:
                content = extractor(str(filepath))
            out_path.parent.mkdir(parents=True, exist_ok=True)
            out_path.write_text(content, encoding='utf-8')
            stats['ok'] += 1
            if i % 50 == 0 or i == len(files):
                print(f"  [{i}/{len(files)}] {stats['ok']} OK, {stats['skip']} skip, {stats['error']} err")
        except Exception as e:
            stats['error'] += 1
            print(f"  ERROR: {filepath}: {e}")

    print(f"\n완료: {stats['ok']} 추출 / {stats['skip']} 스킵 / {stats['error']} 에러")


if __name__ == '__main__':
    main()
